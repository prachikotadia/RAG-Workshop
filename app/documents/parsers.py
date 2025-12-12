"""
Document parsers for extracting text from various file types.

Supports PDF, TXT, MD, and image files (JPG, PNG, GIF, etc.) with metadata extraction.
"""
from pathlib import Path
from typing import Tuple, Dict, Any, Optional
import logging

logger = logging.getLogger(__name__)


def extract_text_from_pdf(path: Path) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text from a PDF file and return (full_text, metadata).
    
    Args:
        path: Path to the PDF file
    
    Returns:
        Tuple of (full_text, metadata_dict)
        Metadata includes: num_pages, source_path, file_type
    """
    try:
        import PyPDF2
    except ImportError:
        logger.error("PyPDF2 not installed. Install with: pip install PyPDF2")
        raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")
    
    try:
        text = ""
        num_pages = 0
        with open(path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        
        metadata = {
            "num_pages": num_pages,
            "source_path": str(path),
            "file_type": "pdf"
        }
        
        return text.strip(), metadata
    except Exception as e:
        logger.error(f"Error parsing PDF {path}: {e}")
        raise ValueError(f"Failed to parse PDF: {str(e)}")


def extract_text_from_text_file(path: Path) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text from a plain text, markdown, JSON, or CSV file and return (full_text, metadata).
    
    Args:
        path: Path to the text/markdown/json/csv file
    
    Returns:
        Tuple of (full_text, metadata_dict)
        Metadata includes: source_path, file_type
    """
    try:
        suffix = path.suffix.lower()
        
        # Handle JSON files - pretty print for readability
        if suffix == ".json":
            try:
                import json
                with open(path, "r", encoding="utf-8") as file:
                    data = json.load(file)
                text = json.dumps(data, indent=2, ensure_ascii=False)
                file_type = "json"
            except json.JSONDecodeError as e:
                logger.warning(f"JSON decode error for {path}: {e} - falling back to raw text")
                with open(path, "r", encoding="utf-8", errors="ignore") as file:
                    text = file.read()
                file_type = "json"
            except Exception as e:
                logger.warning(f"Error parsing JSON {path}: {e} - falling back to raw text")
                with open(path, "r", encoding="utf-8", errors="ignore") as file:
                    text = file.read()
                file_type = "json"
        # Handle CSV files - convert to readable format
        elif suffix == ".csv":
            try:
                import csv
                text_lines = []
                with open(path, "r", encoding="utf-8", errors="ignore") as file:
                    csv_reader = csv.reader(file)
                    for row in csv_reader:
                        text_lines.append(" | ".join(row))
                text = "\n".join(text_lines)
                file_type = "csv"
            except Exception as e:
                logger.warning(f"Error parsing CSV {path}: {e} - falling back to raw text")
                with open(path, "r", encoding="utf-8", errors="ignore") as file:
                    text = file.read()
                file_type = "csv"
        else:
            # Regular text/markdown files
            with open(path, "r", encoding="utf-8", errors="ignore") as file:
                text = file.read()
            
            # Determine file type
            file_type = "markdown" if suffix in [".md", ".markdown"] else "text"
        
        metadata = {
            "source_path": str(path),
            "file_type": file_type
        }
        
        return text, metadata
    except Exception as e:
        logger.error(f"Error parsing text file {path}: {e}")
        raise ValueError(f"Failed to parse text file: {str(e)}")


def extract_text_from_image(path: Path) -> Tuple[str, Dict[str, Any]]:
    """
    Extract metadata from an image file.

    For images, we store metadata (dimensions, format, etc.) as text.
    The actual image analysis will be done by CLIP/BLIP-2 when needed.
    
    Supports: JPG, JPEG, PNG, GIF, WEBP, BMP, HEIC, HEIF, TIFF, TIF, SVG, ICO

    Args:
        path: Path to the image file

    Returns:
        Tuple of (description_text, metadata_dict)
        Description text includes basic image info
        Metadata includes: width, height, format, file_size
    """
    try:
        from PIL import Image
        
        # For HEIC/HEIF files, try to load pillow-heif plugin
        suffix = path.suffix.lower()
        if suffix in [".heic", ".heif"]:
            try:
                from pillow_heif import register_heif_opener
                register_heif_opener()
                logger.debug(f"Registered HEIC/HEIF opener for {path}")
            except ImportError:
                logger.warning(
                    "pillow-heif not installed. HEIC/HEIF support may be limited. "
                    "Install with: pip install pillow-heif"
                )
        
        # Open and process image
        with Image.open(path) as img:
            # Convert to RGB if necessary (for formats like RGBA, P, etc.)
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")
            
            width, height = img.size
            format_name = img.format or "unknown"
            mode = img.mode
            
            # Get file size
            file_size = path.stat().st_size
            
            # Create a basic description
            description = f"Image file: {path.name}\nDimensions: {width}x{height} pixels\nFormat: {format_name}\nColor mode: {mode}\nFile size: {file_size} bytes"
            
            metadata = {
                "width": width,
                "height": height,
                "format": format_name,
                "mode": mode,
                "file_size": file_size,
                "source_path": str(path),
                "file_type": "image"
            }
            
            return description, metadata
            
    except ImportError:
        logger.error("Pillow (PIL) not installed. Install with: pip install pillow")
        raise ImportError("Pillow not installed. Install with: pip install pillow")
    except Exception as e:
        logger.error(f"Error processing image {path}: {e}")
        raise ValueError(f"Failed to process image: {str(e)}")


def extract_text_from_office_document(path: Path) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text from Office documents (.docx, .pptx).
    
    Args:
        path: Path to the Office document
    
    Returns:
        Tuple of (full_text, metadata_dict)
        Metadata includes: source_path, file_type
    """
    suffix = path.suffix.lower()
    file_type = "docx" if suffix == ".docx" else "pptx"
    
    try:
        if suffix == ".docx":
            try:
                from docx import Document
                doc = Document(path)
                text_parts = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)
                text = "\n".join(text_parts)
            except ImportError:
                logger.warning("python-docx not installed. Install with: pip install python-docx - falling back to raw mode")
                # Fallback: try to extract as ZIP (docx is a ZIP file)
                try:
                    import zipfile
                    with zipfile.ZipFile(path, 'r') as zip_ref:
                        # Try to read the main document XML
                        try:
                            xml_content = zip_ref.read('word/document.xml')
                            # Basic text extraction from XML (very simple)
                            import re
                            text = re.sub(r'<[^>]+>', '', xml_content.decode('utf-8', errors='ignore'))
                        except:
                            text = f"Office document: {path.name} (python-docx required for proper parsing)"
                except:
                    text = f"Office document: {path.name} (python-docx required for proper parsing)"
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
            except ImportError:
                logger.warning("python-pptx not installed. Install with: pip install python-pptx - falling back to raw mode")
                # Fallback: try to extract as ZIP (pptx is a ZIP file)
                try:
                    import zipfile
                    with zipfile.ZipFile(path, 'r') as zip_ref:
                        # Try to read slide XML files
                        text_parts = []
                        for name in zip_ref.namelist():
                            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                                try:
                                    xml_content = zip_ref.read(name)
                                    import re
                                    slide_text = re.sub(r'<[^>]+>', '', xml_content.decode('utf-8', errors='ignore'))
                                    if slide_text.strip():
                                        text_parts.append(slide_text)
                                except:
                                    pass
                        text = "\n".join(text_parts) if text_parts else f"Presentation: {path.name} (python-pptx required for proper parsing)"
                except:
                    text = f"Presentation: {path.name} (python-pptx required for proper parsing)"
        else:
            text = f"Office document: {path.name}"
        
        metadata = {
            "source_path": str(path),
            "file_type": file_type
        }
        
        return text, metadata
    except Exception as e:
        logger.error(f"Error parsing Office document {path}: {e}")
        # Fallback to basic text
        return f"Office document: {path.name}\nError: {str(e)}", {
            "source_path": str(path),
            "file_type": file_type,
            "error": str(e)
        }


def extract_text_from_file(path: Path) -> Tuple[str, Dict[str, Any]]:
    """
    Detect file type by suffix and route to the appropriate helper.

    Supported formats:
    - Documents: .pdf, .txt, .md, .markdown, .json, .csv, .docx, .pptx
    - Images: .jpg, .jpeg, .png, .gif, .webp, .bmp, .heic, .heif, .tiff, .tif, .svg, .ico
    
    Images are processed with CLIP (for embeddings) and BLIP-2 (for captioning).

    Args:
        path: Path to the file

    Returns:
        Tuple of (full_text, metadata_dict)

    Raises:
        ValueError: For unsupported file types
    """
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    elif suffix in [".txt", ".text", ".md", ".markdown", ".json", ".csv"]:
        return extract_text_from_text_file(path)
    elif suffix in [".docx", ".pptx"]:
        return extract_text_from_office_document(path)
    elif suffix in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".heic", ".heif", ".tiff", ".tif", ".svg", ".ico"]:
        return extract_text_from_image(path)
    else:
        # Fallback: try to read as text file (never reject)
        logger.warning(f"Unknown file type {suffix} for {path.name} - attempting to read as text file")
        try:
            return extract_text_from_text_file(path)
        except Exception as e:
            raise ValueError(
                f"Unsupported file type: {suffix}. "
                f"Supported: .pdf, .txt, .md, .json, .csv, .docx, .pptx, .jpg, .jpeg, .png, .gif, .webp, .bmp, .heic, .heif, .tiff, .tif, .svg, .ico"
            )

